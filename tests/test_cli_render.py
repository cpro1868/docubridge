from pathlib import Path
import json
import base64

from typer.testing import CliRunner
from docx import Document
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt
from openpyxl import Workbook
from pptx import Presentation

from docubridge.cli import app


runner = CliRunner()

_PNG_1X1 = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+aRX0AAAAASUVORK5CYII="
)


def _write_sample_docx(path: Path) -> None:
    document = Document()
    document.add_paragraph("Sample Title", style="Heading 1")
    document.add_paragraph("Plain paragraph.")
    document.add_paragraph("First bullet", style="List Bullet")
    document.add_paragraph("Second bullet", style="List Bullet")
    path.parent.mkdir(parents=True, exist_ok=True)
    document.save(path)


def _write_rich_text_docx(path: Path) -> None:
    document = Document()
    heading = document.add_paragraph(style="Heading 1")
    heading.add_run("Title ")
    heading.add_run("Bold").bold = True
    heading.add_run(" ")
    heading.add_run("Italic").italic = True
    heading.add_run(" ")
    heading.add_run("Gone").font.strike = True
    heading.add_run(" ")
    _append_hyperlink(heading, "ref", "https://example.com")

    paragraph = document.add_paragraph()
    paragraph.add_run("Body ")
    paragraph.add_run("Bold").bold = True
    paragraph.add_run(" ")
    paragraph.add_run("Italic").italic = True
    paragraph.add_run(" ")
    paragraph.add_run("Gone").font.strike = True
    paragraph.add_run(" ")
    _append_hyperlink(paragraph, "ref", "https://example.com")

    path.parent.mkdir(parents=True, exist_ok=True)
    document.save(path)


def _write_code_text_docx(path: Path) -> None:
    document = Document()
    heading = document.add_paragraph(style="Heading 1")
    heading.add_run("Title ")
    code_run = heading.add_run("code")
    code_run.font.name = "Consolas"

    paragraph = document.add_paragraph()
    paragraph.add_run("Body ")
    body_code = paragraph.add_run("code")
    body_code.font.name = "Consolas"

    table = document.add_table(rows=2, cols=1)
    table.cell(0, 0).text = "Value"
    cell = table.cell(1, 0)
    cell.text = ""
    cell_paragraph = cell.paragraphs[0]
    cell_paragraph.add_run("Cell ")
    cell_code = cell_paragraph.add_run("code")
    cell_code.font.name = "Consolas"

    path.parent.mkdir(parents=True, exist_ok=True)
    document.save(path)


def _write_fragmented_bold_docx(path: Path) -> None:
    document = Document()
    paragraph = document.add_paragraph()
    first = paragraph.add_run("1")
    first.bold = True
    second = paragraph.add_run("、")
    second.bold = True
    third = paragraph.add_run("故事发生地点：")
    third.bold = True
    paragraph.add_run("卡拉尔星东半球西大陆")

    table = document.add_table(rows=2, cols=1)
    table.cell(0, 0).text = "Value"
    cell = table.cell(1, 0)
    cell.text = ""
    cell_paragraph = cell.paragraphs[0]
    cell_first = cell_paragraph.add_run("10 ")
    cell_first.bold = True
    cell_second = cell_paragraph.add_run("天")
    cell_second.bold = True

    paragraph_with_punctuation = document.add_paragraph()
    punct_first = paragraph_with_punctuation.add_run("2")
    punct_first.bold = True
    punct_second = paragraph_with_punctuation.add_run("、")
    punct_second.bold = True
    punct_third = paragraph_with_punctuation.add_run("时代科技水平")
    punct_third.bold = True
    paragraph_with_punctuation.add_run("：整体不如现代地球")

    path.parent.mkdir(parents=True, exist_ok=True)
    document.save(path)


def _write_table_docx(path: Path) -> None:
    document = Document()
    table = document.add_table(rows=3, cols=2)
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Alpha"
    table.cell(1, 1).text = "1"
    table.cell(2, 0).text = "Beta"
    table.cell(2, 1).text = "2"
    path.parent.mkdir(parents=True, exist_ok=True)
    document.save(path)


def _append_hyperlink(paragraph, text: str, href: str) -> None:
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    from docx.opc.constants import RELATIONSHIP_TYPE as RT

    part = paragraph.part
    rel_id = part.relate_to(href, RT.HYPERLINK, is_external=True)
    hyperlink = OxmlElement("w:hyperlink")
    hyperlink.set(qn("r:id"), rel_id)
    run = OxmlElement("w:r")
    text_node = OxmlElement("w:t")
    text_node.text = text
    run.append(text_node)
    hyperlink.append(run)
    paragraph._p.append(hyperlink)


def _get_or_add_paragraph_properties(paragraph):
    properties = paragraph._p.pPr
    if properties is None:
        properties = OxmlElement("w:pPr")
        paragraph._p.insert(0, properties)
    return properties


def _set_paragraph_numbering(paragraph, *, num_id: int, ilvl: int = 0) -> None:
    properties = _get_or_add_paragraph_properties(paragraph)
    num_pr = properties.find(qn("w:numPr"))
    if num_pr is None:
        num_pr = OxmlElement("w:numPr")
        properties.append(num_pr)

    ilvl_element = num_pr.find(qn("w:ilvl"))
    if ilvl_element is None:
        ilvl_element = OxmlElement("w:ilvl")
        num_pr.append(ilvl_element)
    ilvl_element.set(qn("w:val"), str(ilvl))

    num_id_element = num_pr.find(qn("w:numId"))
    if num_id_element is None:
        num_id_element = OxmlElement("w:numId")
        num_pr.append(num_id_element)
    num_id_element.set(qn("w:val"), str(num_id))


def _style_numbering_reference(style) -> tuple[int, int] | None:
    current = style
    while current is not None:
        properties = getattr(current.element, "pPr", None)
        num_pr = getattr(properties, "numPr", None)
        if num_pr is not None and getattr(num_pr, "numId", None) is not None:
            ilvl = getattr(num_pr, "ilvl", None)
            return int(num_pr.numId.val), int(ilvl.val if ilvl is not None else 0)
        current = current.base_style
    return None


def _paragraph_numbering_reference(paragraph) -> tuple[int, int]:
    properties = paragraph._p.pPr
    num_pr = getattr(properties, "numPr", None) if properties is not None else None
    if num_pr is not None and getattr(num_pr, "numId", None) is not None:
        ilvl = getattr(num_pr, "ilvl", None)
        return int(num_pr.numId.val), int(ilvl.val if ilvl is not None else 0)

    reference = _style_numbering_reference(paragraph.style)
    if reference is None:
        raise AssertionError("Paragraph has no numbering reference")
    return reference


def _set_numbering_start_override(document: Document, *, num_id: int, ilvl: int, start_at: int) -> None:
    numbering = document.part.numbering_part.element
    num_nodes = numbering.xpath(f"./w:num[@w:numId='{num_id}']")
    if not num_nodes:
        raise AssertionError(f"Missing numbering definition for numId={num_id}")
    num_node = num_nodes[0]

    override_nodes = num_node.xpath(f"./w:lvlOverride[@w:ilvl='{ilvl}']")
    if override_nodes:
        override_node = override_nodes[0]
    else:
        override_node = OxmlElement("w:lvlOverride")
        override_node.set(qn("w:ilvl"), str(ilvl))
        num_node.append(override_node)

    start_override = override_node.find(qn("w:startOverride"))
    if start_override is None:
        start_override = OxmlElement("w:startOverride")
        override_node.append(start_override)
    start_override.set(qn("w:val"), str(start_at))


def _numbering_start_override(document: Document, *, num_id: int, ilvl: int) -> int | None:
    numbering = document.part.numbering_part.element
    num_nodes = numbering.xpath(f"./w:num[@w:numId='{num_id}']")
    if not num_nodes:
        return None
    override_nodes = num_nodes[0].xpath(f"./w:lvlOverride[@w:ilvl='{ilvl}']")
    if not override_nodes:
        return None
    start_nodes = override_nodes[0].xpath("./w:startOverride/@w:val")
    if not start_nodes:
        return None
    return int(start_nodes[0])


def _write_rich_table_docx(path: Path) -> None:
    document = Document()
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Value"

    first_cell = table.cell(1, 0)
    first_cell.text = ""
    first_paragraph = first_cell.paragraphs[0]
    first_paragraph.add_run("Alpha").bold = True
    first_paragraph.add_run(" ")
    first_paragraph.add_run("Italic").italic = True

    second_cell = table.cell(1, 1)
    second_cell.text = ""
    second_paragraph = second_cell.paragraphs[0]
    second_paragraph.add_run("See ")
    _append_hyperlink(second_paragraph, "ref", "https://example.com")

    path.parent.mkdir(parents=True, exist_ok=True)
    document.save(path)


def _write_table_with_multiline_cell_docx(path: Path) -> None:
    document = Document()
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Alpha"
    table.cell(1, 1).text = "Line one\nLine two"
    path.parent.mkdir(parents=True, exist_ok=True)
    document.save(path)


def _write_image_docx(path: Path, image_path: Path) -> None:
    image_path.parent.mkdir(parents=True, exist_ok=True)
    image_path.write_bytes(_PNG_1X1)
    document = Document()
    document.add_paragraph("Image section")
    document.add_picture(str(image_path))
    path.parent.mkdir(parents=True, exist_ok=True)
    document.save(path)


def _write_template_docx(path: Path) -> None:
    document = Document()
    heading_style = document.styles["Heading 1"]
    heading_style.font.name = "Courier New"
    normal_style = document.styles["Normal"]
    normal_style.font.name = "Arial"
    path.parent.mkdir(parents=True, exist_ok=True)
    document.save(path)


def _write_ordered_docx(path: Path, image_path: Path) -> None:
    image_path.parent.mkdir(parents=True, exist_ok=True)
    image_path.write_bytes(_PNG_1X1)
    document = Document()
    document.add_paragraph("Before image")
    document.add_picture(str(image_path))
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Key"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Order"
    table.cell(1, 1).text = "Preserved"
    document.add_paragraph("After table")
    path.parent.mkdir(parents=True, exist_ok=True)
    document.save(path)


def _write_nested_list_docx(path: Path) -> None:
    document = Document()
    document.add_paragraph("Top", style="List Bullet")
    nested = document.add_paragraph("Nested", style="List Bullet")
    nested.paragraph_format.left_indent = Inches(0.5)
    document.add_paragraph("Top 2", style="List Bullet")
    path.parent.mkdir(parents=True, exist_ok=True)
    document.save(path)


def _write_numbered_list_docx(path: Path) -> None:
    document = Document()
    document.add_paragraph("First", style="List Number")
    document.add_paragraph("Second", style="List Number")
    nested = document.add_paragraph("Nested", style="List Number")
    nested.paragraph_format.left_indent = Inches(0.5)
    document.add_paragraph("Third", style="List Number")
    path.parent.mkdir(parents=True, exist_ok=True)
    document.save(path)


def _write_numbered_start_override_docx(path: Path) -> None:
    document = Document()
    first = document.add_paragraph("Five", style="List Number")
    second = document.add_paragraph("Six", style="List Number")
    num_id, ilvl = _paragraph_numbering_reference(first)
    _set_numbering_start_override(document, num_id=num_id, ilvl=ilvl, start_at=5)
    path.parent.mkdir(parents=True, exist_ok=True)
    document.save(path)


def _write_numbered_heading_docx(path: Path) -> None:
    document = Document()
    reference = document.add_paragraph("Reference", style="List Number")
    num_id, ilvl = _paragraph_numbering_reference(reference)
    heading = document.add_paragraph("Heading With Number", style="Heading 1")
    _set_paragraph_numbering(heading, num_id=num_id, ilvl=ilvl)
    reference._element.getparent().remove(reference._element)
    path.parent.mkdir(parents=True, exist_ok=True)
    document.save(path)


def _write_numbered_headings_with_body_docx(path: Path) -> None:
    document = Document()
    reference = document.add_paragraph("Reference", style="List Number")
    num_id, ilvl = _paragraph_numbering_reference(reference)

    first_heading = document.add_paragraph("First Section", style="Heading 1")
    _set_paragraph_numbering(first_heading, num_id=num_id, ilvl=ilvl)
    document.add_paragraph("Body under first section.")

    second_heading = document.add_paragraph("Second Section", style="Heading 1")
    _set_paragraph_numbering(second_heading, num_id=num_id, ilvl=ilvl)
    document.add_paragraph("Body under second section.")

    reference._element.getparent().remove(reference._element)
    path.parent.mkdir(parents=True, exist_ok=True)
    document.save(path)


def _write_numbered_headings_with_table_docx(path: Path) -> None:
    document = Document()
    reference = document.add_paragraph("Reference", style="List Number")
    num_id, ilvl = _paragraph_numbering_reference(reference)

    first_heading = document.add_paragraph("First Section", style="Heading 1")
    _set_paragraph_numbering(first_heading, num_id=num_id, ilvl=ilvl)

    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Key"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "A"
    table.cell(1, 1).text = "1"

    second_heading = document.add_paragraph("Second Section", style="Heading 1")
    _set_paragraph_numbering(second_heading, num_id=num_id, ilvl=ilvl)

    reference._element.getparent().remove(reference._element)
    path.parent.mkdir(parents=True, exist_ok=True)
    document.save(path)


def _write_release_case_docx(path: Path, image_path: Path) -> None:
    image_path.parent.mkdir(parents=True, exist_ok=True)
    image_path.write_bytes(_PNG_1X1)

    document = Document()
    reference = document.add_paragraph("Reference", style="List Number")
    num_id, ilvl = _paragraph_numbering_reference(reference)
    heading = document.add_paragraph("Heading With Number", style="Heading 1")
    _set_paragraph_numbering(heading, num_id=num_id, ilvl=ilvl)
    reference._element.getparent().remove(reference._element)

    document.add_paragraph("Before image")
    document.add_picture(str(image_path))

    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Key"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Order"
    table.cell(1, 1).text = "Preserved"

    document.add_paragraph("After table")
    path.parent.mkdir(parents=True, exist_ok=True)
    document.save(path)


def _write_sample_xlsx(path: Path) -> None:
    workbook = Workbook()
    sheet1 = workbook.active
    sheet1.title = "Summary"
    sheet1.append(["Name", "Value"])
    sheet1.append(["Alpha", 1])
    sheet1.append(["Beta", 2])
    sheet2 = workbook.create_sheet("Details")
    sheet2.append(["Key", "Text"])
    sheet2.append(["Note", "Hello"])
    path.parent.mkdir(parents=True, exist_ok=True)
    workbook.save(path)


def _write_sample_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Intro"
    slide.placeholders[1].text = "First point\nSecond point"
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)


def _write_image_pptx(path: Path, image_path: Path) -> None:
    image_path.parent.mkdir(parents=True, exist_ok=True)
    image_path.write_bytes(_PNG_1X1)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Images"
    picture = slide.shapes.add_picture(str(image_path), Inches(1), Inches(1))
    picture.name = "deck-image"
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)


def _write_notes_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Notes"
    slide.placeholders[1].text = "Visible point"
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.text = "Speaker note one"
    notes_frame.add_paragraph().text = "Speaker note two"
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)


def _write_table_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Table"
    table_shape = slide.shapes.add_table(3, 2, Inches(1), Inches(1.5), Inches(4), Inches(1.5))
    table = table_shape.table
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Alpha"
    table.cell(1, 1).text = "1"
    table.cell(2, 0).text = "Beta"
    table.cell(2, 1).text = "2"
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)


def _write_rich_table_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Table"
    table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(1.5), Inches(5), Inches(1.2))
    table = table_shape.table
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Value"

    first_cell = table.cell(1, 0)
    first_paragraph = first_cell.text_frame.paragraphs[0]
    first_paragraph.text = ""
    bold_run = first_paragraph.add_run()
    bold_run.text = "Alpha"
    bold_run.font.bold = True

    second_cell = table.cell(1, 1)
    second_paragraph = second_cell.text_frame.paragraphs[0]
    second_paragraph.text = ""
    second_paragraph.add_run().text = "See "
    link_run = second_paragraph.add_run()
    link_run.text = "ref"
    link_run.hyperlink.address = "https://example.com"

    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)


def _write_ordered_pptx(path: Path, image_path: Path) -> None:
    image_path.parent.mkdir(parents=True, exist_ok=True)
    image_path.write_bytes(_PNG_1X1)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Ordered"
    table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(3.4), Inches(4), Inches(1.2))
    table = table_shape.table
    table.cell(0, 0).text = "Key"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Order"
    table.cell(1, 1).text = "Preserved"
    picture = slide.shapes.add_picture(str(image_path), Inches(1), Inches(2.2))
    picture.name = "ordered-image"
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(4), Inches(0.6))
    textbox.text_frame.text = "Before image"
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)


def _write_text_hierarchy_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Hierarchy"
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(5), Inches(2))
    text_frame = textbox.text_frame
    text_frame.text = "Overview paragraph"
    first_item = text_frame.add_paragraph()
    first_item.text = "Top item"
    first_item.level = 0
    nested_item = text_frame.add_paragraph()
    nested_item.text = "Nested item"
    nested_item.level = 1
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)


def _write_rich_text_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Rich Text"
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.4), Inches(6), Inches(1.5))
    text_frame = textbox.text_frame
    paragraph = text_frame.paragraphs[0]
    paragraph.text = ""
    paragraph.add_run().text = "Body "
    bold_run = paragraph.add_run()
    bold_run.text = "Bold"
    bold_run.font.bold = True
    paragraph.add_run().text = " "
    italic_run = paragraph.add_run()
    italic_run.text = "Italic"
    italic_run.font.italic = True
    paragraph.add_run().text = " "
    link_run = paragraph.add_run()
    link_run.text = "ref"
    link_run.hyperlink.address = "https://example.com"
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)


def _write_code_text_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Code"

    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.4), Inches(6), Inches(1.2))
    text_frame = textbox.text_frame
    paragraph = text_frame.paragraphs[0]
    paragraph.text = ""
    paragraph.add_run().text = "Body "
    code_run = paragraph.add_run()
    code_run.text = "code"
    code_run.font.name = "Consolas"

    table_shape = slide.shapes.add_table(2, 1, Inches(1), Inches(3), Inches(4), Inches(1.2))
    table = table_shape.table
    table.cell(0, 0).text = "Value"
    cell_paragraph = table.cell(1, 0).text_frame.paragraphs[0]
    cell_paragraph.text = ""
    cell_paragraph.add_run().text = "Cell "
    cell_code = cell_paragraph.add_run()
    cell_code.text = "code"
    cell_code.font.name = "Consolas"

    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)


def test_parse_command_creates_markdown_from_docx(tmp_path: Path) -> None:
    input_path = tmp_path / "sample.docx"
    output_path = tmp_path / "sample.md"
    _write_sample_docx(input_path)

    result = runner.invoke(
        app,
        [
            "parse",
            str(input_path),
            "-o",
            str(output_path),
        ],
    )

    assert result.exit_code == 0
    assert output_path.exists()
    content = output_path.read_text(encoding="utf-8")
    assert "# Sample Title" in content
    assert "Plain paragraph." in content
    assert "- First bullet" in content
    assert "- Second bullet" in content


def test_parse_command_preserves_rich_text_in_docx_headings_and_paragraphs(tmp_path: Path) -> None:
    input_path = tmp_path / "rich-text.docx"
    output_path = tmp_path / "rich-text.md"
    _write_rich_text_docx(input_path)

    result = runner.invoke(
        app,
        [
            "parse",
            str(input_path),
            "-o",
            str(output_path),
        ],
    )

    assert result.exit_code == 0
    content = output_path.read_text(encoding="utf-8")
    assert "# Title **Bold** *Italic* ~~Gone~~ [ref](https://example.com)" in content
    assert "Body **Bold** *Italic* ~~Gone~~ [ref](https://example.com)" in content


def test_parse_command_preserves_inline_code_in_docx_content(tmp_path: Path) -> None:
    input_path = tmp_path / "code-text.docx"
    output_path = tmp_path / "code-text.md"
    _write_code_text_docx(input_path)

    result = runner.invoke(
        app,
        [
            "parse",
            str(input_path),
            "-o",
            str(output_path),
        ],
    )

    assert result.exit_code == 0
    content = output_path.read_text(encoding="utf-8")
    assert "# Title `code`" in content
    assert "Body `code`" in content
    assert "| Cell `code` |" in content


def test_parse_command_merges_adjacent_bold_runs_into_single_markdown_span(
    tmp_path: Path,
) -> None:
    input_path = tmp_path / "fragmented-bold.docx"
    output_path = tmp_path / "fragmented-bold.md"
    _write_fragmented_bold_docx(input_path)

    result = runner.invoke(
        app,
        [
            "parse",
            str(input_path),
            "-o",
            str(output_path),
        ],
    )

    assert result.exit_code == 0
    content = output_path.read_text(encoding="utf-8")
    assert "**1、故事发生地点：** 卡拉尔星东半球西大陆" in content
    assert "**2、时代科技水平：** 整体不如现代地球" in content
    assert "**1****、****故事发生地点：**" not in content
    assert "| **10 天** |" in content
    assert "**10 ****天**" not in content


def test_parse_command_extracts_docx_tables_as_markdown(tmp_path: Path) -> None:
    input_path = tmp_path / "table.docx"
    output_path = tmp_path / "table.md"
    _write_table_docx(input_path)

    result = runner.invoke(
        app,
        [
            "parse",
            str(input_path),
            "-o",
            str(output_path),
        ],
    )

    assert result.exit_code == 0
    content = output_path.read_text(encoding="utf-8")
    assert "| Name | Value |" in content
    assert "| --- | --- |" in content
    assert "| Alpha | 1 |" in content
    assert "| Beta | 2 |" in content


def test_parse_command_preserves_rich_text_inside_docx_table_cells(tmp_path: Path) -> None:
    input_path = tmp_path / "rich-table.docx"
    output_path = tmp_path / "rich-table.md"
    _write_rich_table_docx(input_path)

    result = runner.invoke(
        app,
        [
            "parse",
            str(input_path),
            "-o",
            str(output_path),
        ],
    )

    assert result.exit_code == 0
    content = output_path.read_text(encoding="utf-8")
    assert "| **Alpha** *Italic* | See [ref](https://example.com) |" in content


def test_parse_command_preserves_multiline_table_cells_as_single_markdown_row(
    tmp_path: Path,
) -> None:
    input_path = tmp_path / "multiline-table.docx"
    output_path = tmp_path / "multiline-table.md"
    _write_table_with_multiline_cell_docx(input_path)

    result = runner.invoke(
        app,
        [
            "parse",
            str(input_path),
            "-o",
            str(output_path),
        ],
    )

    assert result.exit_code == 0
    content = output_path.read_text(encoding="utf-8")
    assert "| Alpha | Line one<br>Line two |" in content
    assert "Line one\nLine two |" not in content


def test_parse_command_exports_docx_images_as_markdown_references(tmp_path: Path) -> None:
    input_path = tmp_path / "image.docx"
    source_image = tmp_path / "source.png"
    output_path = tmp_path / "image.md"
    _write_image_docx(input_path, source_image)

    result = runner.invoke(
        app,
        [
            "parse",
            str(input_path),
            "-o",
            str(output_path),
        ],
    )

    assert result.exit_code == 0
    content = output_path.read_text(encoding="utf-8")
    assert "![source](" in content
    assert "assets/" in content
    assets_dir = output_path.parent / "assets"
    exported = list(assets_dir.iterdir())
    assert exported


def test_parse_command_preserves_docx_block_order_for_images_and_tables(tmp_path: Path) -> None:
    input_path = tmp_path / "ordered.docx"
    source_image = tmp_path / "ordered.png"
    output_path = tmp_path / "ordered.md"
    _write_ordered_docx(input_path, source_image)

    result = runner.invoke(
        app,
        [
            "parse",
            str(input_path),
            "-o",
            str(output_path),
        ],
    )

    assert result.exit_code == 0
    content = output_path.read_text(encoding="utf-8")
    before_index = content.index("Before image")
    image_index = content.index("![ordered](")
    table_index = content.index("| Key | Value |")
    after_index = content.index("After table")
    assert before_index < image_index < table_index < after_index


def test_parse_command_preserves_simple_nested_bullet_lists(tmp_path: Path) -> None:
    input_path = tmp_path / "nested.docx"
    output_path = tmp_path / "nested.md"
    _write_nested_list_docx(input_path)

    result = runner.invoke(
        app,
        [
            "parse",
            str(input_path),
            "-o",
            str(output_path),
        ],
    )

    assert result.exit_code == 0
    content = output_path.read_text(encoding="utf-8")
    assert "- Top" in content
    assert "  - Nested" in content
    assert "- Top 2" in content


def test_parse_command_preserves_simple_numbered_list_progression(tmp_path: Path) -> None:
    input_path = tmp_path / "numbered.docx"
    output_path = tmp_path / "numbered.md"
    _write_numbered_list_docx(input_path)

    result = runner.invoke(
        app,
        [
            "parse",
            str(input_path),
            "-o",
            str(output_path),
        ],
    )

    assert result.exit_code == 0
    content = output_path.read_text(encoding="utf-8")
    assert "1. First" in content
    assert "2. Second" in content
    assert "  1. Nested" in content
    assert "3. Third" in content


def test_parse_command_preserves_original_docx_numbering_values(tmp_path: Path) -> None:
    input_path = tmp_path / "numbered-start.docx"
    output_path = tmp_path / "numbered-start.md"
    _write_numbered_start_override_docx(input_path)

    result = runner.invoke(
        app,
        [
            "parse",
            str(input_path),
            "-o",
            str(output_path),
        ],
    )

    assert result.exit_code == 0
    content = output_path.read_text(encoding="utf-8")
    assert "5. Five" in content
    assert "6. Six" in content


def test_parse_command_preserves_numbered_docx_headings(tmp_path: Path) -> None:
    input_path = tmp_path / "heading-numbered.docx"
    output_path = tmp_path / "heading-numbered.md"
    _write_numbered_heading_docx(input_path)

    result = runner.invoke(
        app,
        [
            "parse",
            str(input_path),
            "-o",
            str(output_path),
        ],
    )

    assert result.exit_code == 0
    content = output_path.read_text(encoding="utf-8")
    assert "# 1. Heading With Number" in content


def test_parse_command_keeps_heading_number_progression_across_body_paragraphs(
    tmp_path: Path,
) -> None:
    input_path = tmp_path / "heading-numbered-body.docx"
    output_path = tmp_path / "heading-numbered-body.md"
    _write_numbered_headings_with_body_docx(input_path)

    result = runner.invoke(
        app,
        [
            "parse",
            str(input_path),
            "-o",
            str(output_path),
        ],
    )

    assert result.exit_code == 0
    content = output_path.read_text(encoding="utf-8")
    assert "# 1. First Section" in content
    assert "# 2. Second Section" in content


def test_parse_command_keeps_heading_number_progression_across_tables(
    tmp_path: Path,
) -> None:
    input_path = tmp_path / "heading-numbered-table.docx"
    output_path = tmp_path / "heading-numbered-table.md"
    _write_numbered_headings_with_table_docx(input_path)

    result = runner.invoke(
        app,
        [
            "parse",
            str(input_path),
            "-o",
            str(output_path),
        ],
    )

    assert result.exit_code == 0
    content = output_path.read_text(encoding="utf-8")
    assert "# 1. First Section" in content
    assert "# 2. Second Section" in content


def test_parse_command_release_case_preserves_docx_numbering_assets_and_order(
    tmp_path: Path,
) -> None:
    input_path = tmp_path / "release.docx"
    source_image = tmp_path / "release.png"
    output_path = tmp_path / "release.md"
    _write_release_case_docx(input_path, source_image)

    result = runner.invoke(
        app,
        [
            "parse",
            str(input_path),
            "-o",
            str(output_path),
        ],
    )

    assert result.exit_code == 0
    content = output_path.read_text(encoding="utf-8")
    heading_index = content.index("# 1. Heading With Number")
    before_index = content.index("Before image")
    image_index = content.index("![release](")
    table_index = content.index("| Key | Value |")
    after_index = content.index("After table")
    assert heading_index < before_index < image_index < table_index < after_index
    assert "assets/" in content
    assets_dir = output_path.parent / "assets"
    assert list(assets_dir.iterdir())


def test_parse_help_includes_json_option_and_description() -> None:
    result = runner.invoke(app, ["parse", "--help"])
    normalized = " ".join(result.stdout.split())

    assert result.exit_code == 0
    assert "Parse supported documents to Markdown." in normalized
    assert "Currently supports .docx, .xlsx, and .pptx input." in normalized
    assert "--json" in result.stdout
    assert "--output" in result.stdout


def test_parse_command_supports_json_output(tmp_path: Path) -> None:
    input_path = tmp_path / "sample.docx"
    output_path = tmp_path / "sample.md"
    _write_sample_docx(input_path)

    result = runner.invoke(
        app,
        [
            "parse",
            str(input_path),
            "-o",
            str(output_path),
            "--json",
        ],
    )

    assert result.exit_code == 0
    payload = json.loads(result.stdout)
    assert payload == {
        "success": True,
        "output_path": output_path.as_posix(),
        "diagnostics": [],
    }


def test_parse_missing_input_supports_json_output(tmp_path: Path) -> None:
    input_path = tmp_path / "missing.docx"
    output_path = tmp_path / "sample.md"

    result = runner.invoke(
        app,
        [
            "parse",
            str(input_path),
            "-o",
            str(output_path),
            "--json",
        ],
    )

    assert result.exit_code == 5
    payload = json.loads(result.stdout)
    assert payload["success"] is False
    assert payload["output_path"] == output_path.as_posix()
    assert payload["diagnostics"][0]["code"] == "INPUT_FILE_NOT_FOUND"


def test_parse_command_creates_markdown_from_xlsx(tmp_path: Path) -> None:
    input_path = tmp_path / "sample.xlsx"
    output_path = tmp_path / "sample.md"
    _write_sample_xlsx(input_path)

    result = runner.invoke(
        app,
        [
            "parse",
            str(input_path),
            "-o",
            str(output_path),
        ],
    )

    assert result.exit_code == 0
    content = output_path.read_text(encoding="utf-8")
    assert "## Summary" in content
    assert "| Name | Value |" in content
    assert "## Details" in content


def test_parse_command_creates_markdown_from_pptx(tmp_path: Path) -> None:
    input_path = tmp_path / "sample.pptx"
    output_path = tmp_path / "sample.md"
    _write_sample_pptx(input_path)

    result = runner.invoke(
        app,
        [
            "parse",
            str(input_path),
            "-o",
            str(output_path),
        ],
    )

    assert result.exit_code == 0
    content = output_path.read_text(encoding="utf-8")
    assert "### Slide 1" in content
    assert "#### Intro" in content
    assert "- First point" in content


def test_parse_command_exports_pptx_images_as_markdown_references(tmp_path: Path) -> None:
    input_path = tmp_path / "images.pptx"
    source_image = tmp_path / "diagram.png"
    output_path = tmp_path / "images.md"
    _write_image_pptx(input_path, source_image)

    result = runner.invoke(
        app,
        [
            "parse",
            str(input_path),
            "-o",
            str(output_path),
        ],
    )

    assert result.exit_code == 0
    content = output_path.read_text(encoding="utf-8")
    assert "![deck-image](" in content
    assert "assets/" in content
    media_dir = output_path.parent / "assets"
    exported = list(media_dir.iterdir())
    assert exported


def test_parse_command_extracts_pptx_notes_as_blockquotes(tmp_path: Path) -> None:
    input_path = tmp_path / "notes.pptx"
    output_path = tmp_path / "notes.md"
    _write_notes_pptx(input_path)

    result = runner.invoke(
        app,
        [
            "parse",
            str(input_path),
            "-o",
            str(output_path),
        ],
    )

    assert result.exit_code == 0
    content = output_path.read_text(encoding="utf-8")
    assert "> Notes:" in content
    assert "> Speaker note one" in content
    assert "> Speaker note two" in content


def test_parse_command_extracts_pptx_tables_as_markdown(tmp_path: Path) -> None:
    input_path = tmp_path / "table.pptx"
    output_path = tmp_path / "table.md"
    _write_table_pptx(input_path)

    result = runner.invoke(
        app,
        [
            "parse",
            str(input_path),
            "-o",
            str(output_path),
        ],
    )

    assert result.exit_code == 0
    content = output_path.read_text(encoding="utf-8")
    assert "| Name | Value |" in content
    assert "| --- | --- |" in content
    assert "| Alpha | 1 |" in content
    assert "| Beta | 2 |" in content


def test_parse_command_preserves_rich_text_inside_pptx_table_cells(tmp_path: Path) -> None:
    input_path = tmp_path / "rich-table.pptx"
    output_path = tmp_path / "rich-table.md"
    _write_rich_table_pptx(input_path)

    result = runner.invoke(
        app,
        [
            "parse",
            str(input_path),
            "-o",
            str(output_path),
        ],
    )

    assert result.exit_code == 0
    content = output_path.read_text(encoding="utf-8")
    assert "| **Alpha** | See [ref](https://example.com) |" in content


def test_parse_command_preserves_pptx_block_order_for_text_images_and_tables(tmp_path: Path) -> None:
    input_path = tmp_path / "ordered.pptx"
    source_image = tmp_path / "ordered.png"
    output_path = tmp_path / "ordered.md"
    _write_ordered_pptx(input_path, source_image)

    result = runner.invoke(
        app,
        [
            "parse",
            str(input_path),
            "-o",
            str(output_path),
        ],
    )

    assert result.exit_code == 0
    content = output_path.read_text(encoding="utf-8")
    text_index = content.index("Before image")
    image_index = content.index("![ordered-image](")
    table_index = content.index("| Key | Value |")
    assert text_index < image_index < table_index


def test_parse_command_distinguishes_pptx_plain_text_and_nested_list_items(tmp_path: Path) -> None:
    input_path = tmp_path / "hierarchy.pptx"
    output_path = tmp_path / "hierarchy.md"
    _write_text_hierarchy_pptx(input_path)

    result = runner.invoke(
        app,
        [
            "parse",
            str(input_path),
            "-o",
            str(output_path),
        ],
    )

    assert result.exit_code == 0
    content = output_path.read_text(encoding="utf-8")
    assert "Overview paragraph" in content
    assert "- Top item" in content
    assert "  - Nested item" in content


def test_parse_command_preserves_rich_text_in_pptx_textboxes(tmp_path: Path) -> None:
    input_path = tmp_path / "rich-text.pptx"
    output_path = tmp_path / "rich-text.md"
    _write_rich_text_pptx(input_path)

    result = runner.invoke(
        app,
        [
            "parse",
            str(input_path),
            "-o",
            str(output_path),
        ],
    )

    assert result.exit_code == 0
    content = output_path.read_text(encoding="utf-8")
    assert "#### Rich Text" in content
    assert "Body **Bold** *Italic* [ref](https://example.com)" in content


def test_parse_command_preserves_inline_code_in_pptx_content(tmp_path: Path) -> None:
    input_path = tmp_path / "code-text.pptx"
    output_path = tmp_path / "code-text.md"
    _write_code_text_pptx(input_path)

    result = runner.invoke(
        app,
        [
            "parse",
            str(input_path),
            "-o",
            str(output_path),
        ],
    )

    assert result.exit_code == 0
    content = output_path.read_text(encoding="utf-8")
    assert "Body `code`" in content
    assert "| Cell `code` |" in content


def test_render_command_creates_docx(tmp_path: Path) -> None:
    output = tmp_path / "out.docx"
    result = runner.invoke(
        app,
        [
            "render",
            "tests/fixtures/sample.md",
            "-o",
            str(output),
            "--style",
            "tests/fixtures/style.yaml",
        ],
    )

    assert result.exit_code == 0
    assert output.exists()


def test_render_help_includes_json_option_and_description() -> None:
    result = runner.invoke(app, ["render", "--help"])

    assert result.exit_code == 0
    assert "Render Markdown to .docx." in result.stdout
    assert "--json" in result.stdout
    assert "--style" in result.stdout
    assert "--template" in result.stdout


def test_render_help_describes_style_and_template_roles() -> None:
    result = runner.invoke(app, ["render", "--help"])

    assert result.exit_code == 0
    assert "YAML style profile" in result.stdout
    assert "explicit" in result.stdout
    assert "formatting rules" in result.stdout
    assert "Use a .docx template as the host document." in result.stdout
    assert "YAML still defines explicit style intent." in result.stdout


def test_render_returns_exit_code_4_for_missing_template_file(tmp_path: Path) -> None:
    missing_template = tmp_path / "missing-template.docx"
    output = tmp_path / "out.docx"

    result = runner.invoke(
        app,
        [
            "render",
            "tests/fixtures/sample.md",
            "-o",
            str(output),
            "--style",
            "tests/fixtures/style.yaml",
            "--template",
            str(missing_template),
        ],
    )

    assert result.exit_code == 4
    assert "template" in result.output.lower()


def test_render_missing_template_supports_json_output(tmp_path: Path) -> None:
    missing_template = tmp_path / "missing-template.docx"
    output = tmp_path / "out.docx"

    result = runner.invoke(
        app,
        [
            "render",
            "tests/fixtures/sample.md",
            "-o",
            str(output),
            "--style",
            "tests/fixtures/style.yaml",
            "--template",
            str(missing_template),
            "--json",
        ],
    )

    assert result.exit_code == 4
    payload = json.loads(result.stdout)
    assert payload["success"] is False
    assert payload["output_path"] == output.as_posix()
    assert payload["diagnostics"][0]["code"] == "TEMPLATE_VALIDATION_ERROR"


def test_render_returns_exit_code_4_for_missing_template_numbering_resource(tmp_path: Path) -> None:
    input_path = tmp_path / "ordered.md"
    input_path.write_text("1. First\n", encoding="utf-8")
    style_path = tmp_path / "ordered-style.yaml"
    style_path.write_text(
        "\n".join(
            [
                "defaults:",
                "  font_size: 12",
                "elements:",
                "  paragraph:",
                "    font_size: 12",
                "  ordered_list:",
                "    template_style: Normal",
                "    numbering_style: Normal",
            ]
        ),
        encoding="utf-8",
    )
    template_path = tmp_path / "template.docx"
    _write_template_docx(template_path)
    output = tmp_path / "out.docx"

    result = runner.invoke(
        app,
        [
            "render",
            str(input_path),
            "-o",
            str(output),
            "--style",
            str(style_path),
            "--template",
            str(template_path),
        ],
    )

    assert result.exit_code == 4
    assert "numbering" in result.output.lower()


def test_render_command_uses_template_file_for_output_styles(tmp_path: Path) -> None:
    input_path = tmp_path / "template.md"
    input_path.write_text("# Template Title\n\nBody paragraph.\n", encoding="utf-8")

    style_path = tmp_path / "style.yaml"
    style_path.write_text(
        "\n".join(
            [
                "defaults:",
                "  font_name: Times New Roman",
                "  font_size: 12",
                "elements:",
                "  paragraph:",
                "    font_size: 12",
                "  heading1:",
                "    font_size: 18",
            ]
        ),
        encoding="utf-8",
    )

    template_path = tmp_path / "template.docx"
    _write_template_docx(template_path)
    output = tmp_path / "out.docx"

    result = runner.invoke(
        app,
        [
            "render",
            str(input_path),
            "-o",
            str(output),
            "--style",
            str(style_path),
            "--template",
            str(template_path),
        ],
    )

    assert result.exit_code == 0
    document = Document(output)
    assert document.styles["Heading 1"].font.name == "Courier New"
    assert document.styles["Normal"].font.name == "Arial"


def test_render_command_allows_builtin_template_styles_without_template(tmp_path: Path) -> None:
    input_path = tmp_path / "builtin-style.md"
    input_path.write_text("# Title\n\nBody paragraph.\n", encoding="utf-8")

    style_path = tmp_path / "builtin-style.yaml"
    style_path.write_text(
        "\n".join(
            [
                "defaults:",
                "  font_name: Times New Roman",
                "  font_size: 12",
                "elements:",
                "  paragraph:",
                "    template_style: Normal",
                "    font_size: 12",
                "  heading1:",
                "    template_style: Heading 1",
                "    font_size: 18",
            ]
        ),
        encoding="utf-8",
    )
    output = tmp_path / "out.docx"

    result = runner.invoke(
        app,
        [
            "render",
            str(input_path),
            "-o",
            str(output),
            "--style",
            str(style_path),
        ],
    )

    assert result.exit_code == 0
    assert output.exists()
    document = Document(output)
    assert document.paragraphs[0].text == "Title"
    assert document.paragraphs[0].style.name == "Heading 1"
    assert document.paragraphs[1].text == "Body paragraph."
    assert document.paragraphs[1].style.name == "Normal"


def test_render_command_applies_block_styles_from_template_backed_yaml(tmp_path: Path) -> None:
    input_path = tmp_path / "structured.md"
    input_path.write_text(
        "\n".join(
            [
                "# Styled Title",
                "",
                "1. First",
                "",
                "- Bullet",
                "",
                "> Quoted",
                "",
                "| Name | Value |",
                "| --- | --- |",
                "| Alpha | 1 |",
                "",
                "```python",
                "print(1)",
                "```",
            ]
        ),
        encoding="utf-8",
    )

    style_path = tmp_path / "structured.yaml"
    style_path.write_text(
        "\n".join(
            [
                "defaults:",
                "  font_name: Times New Roman",
                "  font_size: 12",
                "elements:",
                "  paragraph:",
                "    font_size: 12",
                "  heading1:",
                "    font_size: 18",
                "  ordered_list:",
                "    template_style: List Number",
                "  unordered_list:",
                "    template_style: List Bullet",
                "  quote:",
                "    template_style: Quote",
                "  table:",
                "    template_style: Table Grid",
                "  code_block:",
                "    template_style: Quote",
            ]
        ),
        encoding="utf-8",
    )

    template_path = tmp_path / "template.docx"
    _write_template_docx(template_path)
    output = tmp_path / "out.docx"

    result = runner.invoke(
        app,
        [
            "render",
            str(input_path),
            "-o",
            str(output),
            "--style",
            str(style_path),
            "--template",
            str(template_path),
        ],
    )

    assert result.exit_code == 0

    document = Document(output)
    paragraph_styles = {paragraph.text: paragraph.style.name for paragraph in document.paragraphs}
    assert paragraph_styles["Styled Title"] == "Heading 1"
    assert paragraph_styles["First"] == "List Number"
    assert paragraph_styles["Bullet"] == "List Bullet"
    assert paragraph_styles["> Quoted"] == "Quote"
    bullet_paragraph = next(paragraph for paragraph in document.paragraphs if paragraph.text == "Bullet")
    assert _paragraph_numbering_reference(bullet_paragraph) is not None
    assert document.tables[0].style.name == "Table Grid"
    assert document.tables[1].cell(0, 0).paragraphs[0].style.name == "Quote"


def test_render_command_binds_nested_ordered_list_to_nested_numbering_level(tmp_path: Path) -> None:
    input_path = tmp_path / "nested-ordered.md"
    input_path.write_text(
        "\n".join(
            [
                "1. Parent",
                "   1. Child",
            ]
        ),
        encoding="utf-8",
    )

    style_path = tmp_path / "nested-style.yaml"
    style_path.write_text(
        "\n".join(
            [
                "elements:",
                "  paragraph:",
                "    template_style: Normal",
                "  ordered_list:",
                "    template_style: List Number",
            ]
        ),
        encoding="utf-8",
    )

    template_path = tmp_path / "template.docx"
    _write_template_docx(template_path)
    output = tmp_path / "out.docx"

    result = runner.invoke(
        app,
        [
            "render",
            str(input_path),
            "-o",
            str(output),
            "--style",
            str(style_path),
            "--template",
            str(template_path),
        ],
    )

    assert result.exit_code == 0
    document = Document(output)
    assert [paragraph.text for paragraph in document.paragraphs] == ["Parent", "Child"]
    assert _paragraph_numbering_reference(document.paragraphs[0])[1] == 0
    assert _paragraph_numbering_reference(document.paragraphs[1])[1] == 1


def test_render_command_restarts_ordered_list_with_markdown_start_value(tmp_path: Path) -> None:
    input_path = tmp_path / "ordered-restart.md"
    input_path.write_text(
        "\n".join(
            [
                "1. First",
                "2. Second",
                "",
                "Break paragraph.",
                "",
                "3. Third",
                "4. Fourth",
            ]
        ),
        encoding="utf-8",
    )

    style_path = tmp_path / "ordered-style.yaml"
    style_path.write_text(
        "\n".join(
            [
                "elements:",
                "  paragraph:",
                "    template_style: Normal",
                "  ordered_list:",
                "    template_style: List Number",
            ]
        ),
        encoding="utf-8",
    )

    template_path = tmp_path / "template.docx"
    _write_template_docx(template_path)
    output = tmp_path / "out.docx"

    result = runner.invoke(
        app,
        [
            "render",
            str(input_path),
            "-o",
            str(output),
            "--style",
            str(style_path),
            "--template",
            str(template_path),
        ],
    )

    assert result.exit_code == 0
    document = Document(output)
    ordered_paragraphs = [paragraph for paragraph in document.paragraphs if paragraph.text in {"First", "Second", "Third", "Fourth"}]
    first_ref = _paragraph_numbering_reference(ordered_paragraphs[0])
    second_ref = _paragraph_numbering_reference(ordered_paragraphs[1])
    third_ref = _paragraph_numbering_reference(ordered_paragraphs[2])
    fourth_ref = _paragraph_numbering_reference(ordered_paragraphs[3])
    assert first_ref[0] == second_ref[0]
    assert third_ref[0] == fourth_ref[0]
    assert third_ref[0] != first_ref[0]
    assert _numbering_start_override(document, num_id=third_ref[0], ilvl=0) == 3


def test_render_command_tracks_ordered_sequences_independently_per_level(tmp_path: Path) -> None:
    input_path = tmp_path / "nested-ordered-start.md"
    input_path.write_text(
        "\n".join(
            [
                "1. Parent One",
                "   1. Child A",
                "   2. Child B",
                "2. Parent Two",
                "   1. Child C",
                "   2. Child D",
            ]
        ),
        encoding="utf-8",
    )

    style_path = tmp_path / "ordered-style.yaml"
    style_path.write_text(
        "\n".join(
            [
                "elements:",
                "  paragraph:",
                "    template_style: Normal",
                "  ordered_list:",
                "    template_style: List Number",
            ]
        ),
        encoding="utf-8",
    )

    template_path = tmp_path / "template.docx"
    _write_template_docx(template_path)
    output = tmp_path / "out.docx"

    result = runner.invoke(
        app,
        [
            "render",
            str(input_path),
            "-o",
            str(output),
            "--style",
            str(style_path),
            "--template",
            str(template_path),
        ],
    )

    assert result.exit_code == 0
    document = Document(output)
    ordered_paragraphs = [
        paragraph for paragraph in document.paragraphs
        if paragraph.text in {"Parent One", "Child A", "Child B", "Parent Two", "Child C", "Child D"}
    ]
    parent_one_ref = _paragraph_numbering_reference(ordered_paragraphs[0])
    child_a_ref = _paragraph_numbering_reference(ordered_paragraphs[1])
    child_b_ref = _paragraph_numbering_reference(ordered_paragraphs[2])
    parent_two_ref = _paragraph_numbering_reference(ordered_paragraphs[3])
    child_c_ref = _paragraph_numbering_reference(ordered_paragraphs[4])
    child_d_ref = _paragraph_numbering_reference(ordered_paragraphs[5])
    assert parent_one_ref is not None and child_a_ref is not None and child_b_ref is not None
    assert parent_two_ref is not None and child_c_ref is not None and child_d_ref is not None
    assert parent_one_ref[1] == 0
    assert parent_two_ref[1] == 0
    assert child_a_ref[1] == 1
    assert child_b_ref[1] == 1
    assert child_c_ref[1] == 1
    assert child_d_ref[1] == 1
    assert parent_one_ref[0] == parent_two_ref[0]
    assert child_a_ref[0] == child_b_ref[0]
    assert child_c_ref[0] == child_d_ref[0]
    assert child_c_ref[0] != child_a_ref[0]


def test_render_command_binds_nested_unordered_list_to_nested_numbering_level(tmp_path: Path) -> None:
    input_path = tmp_path / "nested-unordered.md"
    input_path.write_text(
        "\n".join(
            [
                "- Parent",
                "  - Child",
            ]
        ),
        encoding="utf-8",
    )

    style_path = tmp_path / "nested-style.yaml"
    style_path.write_text(
        "\n".join(
            [
                "elements:",
                "  paragraph:",
                "    template_style: Normal",
                "  unordered_list:",
                "    template_style: List Bullet",
            ]
        ),
        encoding="utf-8",
    )

    template_path = tmp_path / "template.docx"
    _write_template_docx(template_path)
    output = tmp_path / "out.docx"

    result = runner.invoke(
        app,
        [
            "render",
            str(input_path),
            "-o",
            str(output),
            "--style",
            str(style_path),
            "--template",
            str(template_path),
        ],
    )

    assert result.exit_code == 0
    document = Document(output)
    assert [paragraph.text for paragraph in document.paragraphs] == ["Parent", "Child"]
    assert _paragraph_numbering_reference(document.paragraphs[0])[1] == 0
    assert _paragraph_numbering_reference(document.paragraphs[1])[1] == 1


def test_render_command_supports_json_output(tmp_path: Path) -> None:
    output = tmp_path / "out.docx"
    result = runner.invoke(
        app,
        [
            "render",
            "tests/fixtures/sample.md",
            "-o",
            str(output),
            "--style",
            "tests/fixtures/style.yaml",
            "--json",
        ],
    )

    assert result.exit_code == 0
    payload = json.loads(result.stdout)
    assert payload == {
        "success": True,
        "output_path": output.as_posix(),
        "diagnostics": [],
    }
    assert output.exists()


def test_style_list_prints_builtin_profiles() -> None:
    result = runner.invoke(app, ["style", "list"])

    assert result.exit_code == 0
    assert "academic" in result.stdout
    assert "business" in result.stdout


def test_render_returns_exit_code_4_for_invalid_style(tmp_path: Path) -> None:
    bad_style = tmp_path / "bad.yaml"
    bad_style.write_text("defaults: []\n", encoding="utf-8")
    output = tmp_path / "out.docx"

    result = runner.invoke(
        app,
        [
            "render",
            "tests/fixtures/sample.md",
            "-o",
            str(output),
            "--style",
            str(bad_style),
        ],
    )

    assert result.exit_code == 4
    assert "dictionary" in result.output.lower()


def test_render_invalid_style_supports_json_output(tmp_path: Path) -> None:
    bad_style = tmp_path / "bad.yaml"
    bad_style.write_text("defaults: []\n", encoding="utf-8")
    output = tmp_path / "out.docx"

    result = runner.invoke(
        app,
        [
            "render",
            "tests/fixtures/sample.md",
            "-o",
            str(output),
            "--style",
            str(bad_style),
            "--json",
        ],
    )

    assert result.exit_code == 4
    payload = json.loads(result.stdout)
    assert payload["success"] is False
    assert payload["output_path"] == output.as_posix()
    assert payload["diagnostics"]
    assert payload["diagnostics"][0]["code"] == "STYLE_VALIDATION_ERROR"


def test_render_returns_exit_code_5_for_missing_markdown_file(tmp_path: Path) -> None:
    missing_input = tmp_path / "missing.md"
    output = tmp_path / "out.docx"

    result = runner.invoke(
        app,
        [
            "render",
            str(missing_input),
            "-o",
            str(output),
            "--style",
            "tests/fixtures/style.yaml",
        ],
    )

    assert result.exit_code == 5
    assert "not found" in result.output.lower() or "no such file" in result.output.lower()


def test_render_missing_markdown_file_supports_json_output(tmp_path: Path) -> None:
    missing_input = tmp_path / "missing.md"
    output = tmp_path / "out.docx"

    result = runner.invoke(
        app,
        [
            "render",
            str(missing_input),
            "-o",
            str(output),
            "--style",
            "tests/fixtures/style.yaml",
            "--json",
        ],
    )

    assert result.exit_code == 5
    payload = json.loads(result.stdout)
    assert payload["success"] is False
    assert payload["output_path"] == output.as_posix()
    assert payload["diagnostics"][0]["code"] == "INPUT_FILE_NOT_FOUND"


def test_render_returns_exit_code_5_for_missing_style_file(tmp_path: Path) -> None:
    missing_style = tmp_path / "missing-style.yaml"
    output = tmp_path / "out.docx"

    result = runner.invoke(
        app,
        [
            "render",
            "tests/fixtures/sample.md",
            "-o",
            str(output),
            "--style",
            str(missing_style),
        ],
    )

    assert result.exit_code == 5
    assert "not found" in result.output.lower() or "no such file" in result.output.lower()


def test_render_accepts_heading2_only_documents_when_profile_defines_heading2(
    tmp_path: Path,
) -> None:
    input_path = tmp_path / "heading2-only.md"
    input_path.write_text("## Second level title\n", encoding="utf-8")

    style_path = tmp_path / "heading2-only.yaml"
    style_path.write_text(
        "\n".join(
            [
                "defaults:",
                "  font_name: Times New Roman",
                "elements:",
                "  paragraph:",
                "    based_on: Normal",
                "    font_size: 12",
                "  heading2:",
                "    based_on: Normal",
                "    font_size: 16",
            ]
        ),
        encoding="utf-8",
    )
    output = tmp_path / "heading2-only.docx"

    result = runner.invoke(
        app,
        [
            "render",
            str(input_path),
            "-o",
            str(output),
            "--style",
            str(style_path),
        ],
    )

    assert result.exit_code == 0
    assert output.exists()


def test_render_returns_exit_code_4_for_style_root_that_is_not_a_mapping(
    tmp_path: Path,
) -> None:
    bad_style = tmp_path / "bad-root.yaml"
    bad_style.write_text("[]\n", encoding="utf-8")
    output = tmp_path / "out.docx"

    result = runner.invoke(
        app,
        [
            "render",
            "tests/fixtures/sample.md",
            "-o",
            str(output),
            "--style",
            str(bad_style),
        ],
    )

    assert result.exit_code == 4
    assert "mapping" in result.output.lower()


def test_render_returns_exit_code_5_for_non_utf8_markdown_input(
    tmp_path: Path,
) -> None:
    bad_input = tmp_path / "bad-input.md"
    bad_input.write_bytes(b"\xff\xfe\xff")
    output = tmp_path / "out.docx"

    result = runner.invoke(
        app,
        [
            "render",
            str(bad_input),
            "-o",
            str(output),
            "--style",
            "tests/fixtures/style.yaml",
        ],
    )

    assert result.exit_code == 5
    assert "utf-8" in result.output.lower() or "decode" in result.output.lower()


def test_doctor_without_inputs_reports_environment_ok() -> None:
    result = runner.invoke(app, ["doctor"])

    assert result.exit_code == 0
    assert "Environment OK" in result.stdout
    assert "environment: ok" in result.stdout.lower()


def test_doctor_help_includes_json_option_and_description() -> None:
    result = runner.invoke(app, ["doctor", "--help"])

    assert result.exit_code == 0
    assert "Run environment and task-level preflight checks." in result.stdout
    assert "--json" in result.stdout
    assert "--style" in result.stdout
    assert "--template" in result.stdout


def test_doctor_without_inputs_supports_json_output() -> None:
    result = runner.invoke(app, ["doctor", "--json"])

    assert result.exit_code == 0
    payload = json.loads(result.stdout)
    assert payload["success"] is True
    assert payload["summary"] == "Environment OK"
    assert payload["checks"] == [{"name": "environment", "status": "ok"}]
    assert payload["warnings"] == []


def test_doctor_with_input_and_style_reports_task_ok() -> None:
    result = runner.invoke(
        app,
        [
            "doctor",
            "tests/fixtures/sample.md",
            "--style",
            "tests/fixtures/style.yaml",
        ],
    )

    assert result.exit_code == 0
    assert "Task OK" in result.stdout
    assert "markdown:" in result.stdout.lower()
    assert "style:" in result.stdout.lower()
    assert "style-resolution:" in result.stdout.lower()
    assert "environment:" in result.stdout.lower()


def test_doctor_with_input_and_style_supports_json_output() -> None:
    result = runner.invoke(
        app,
        [
            "doctor",
            "tests/fixtures/sample.md",
            "--style",
            "tests/fixtures/style.yaml",
            "--json",
        ],
    )

    assert result.exit_code == 0
    payload = json.loads(result.stdout)
    assert payload["success"] is True
    assert payload["summary"] == "Task OK"
    assert payload["checks"] == [
        {"name": "environment", "status": "ok"},
        {"name": "markdown", "status": "ok", "detail": "tests/fixtures/sample.md"},
        {"name": "style", "status": "ok", "detail": "tests/fixtures/style.yaml"},
        {"name": "style-resolution", "status": "ok"},
    ]
    assert payload["warnings"] == [
        "markdown contains images that may be degraded or replaced"
    ]


def test_doctor_accepts_heading3_via_heading1_fallback(tmp_path: Path) -> None:
    input_path = tmp_path / "heading3-fallback.md"
    input_path.write_text("### Third level title\n", encoding="utf-8")

    result = runner.invoke(
        app,
        [
            "doctor",
            str(input_path),
            "--style",
            "tests/fixtures/style.yaml",
        ],
    )

    assert result.exit_code == 0
    assert "style-resolution: ok" in result.stdout.lower()
    assert "task ok" in result.stdout.lower()


def test_doctor_returns_exit_code_5_for_missing_input(tmp_path: Path) -> None:
    result = runner.invoke(
        app,
        [
            "doctor",
            str(tmp_path / "missing.md"),
            "--style",
            "tests/fixtures/style.yaml",
        ],
    )

    assert result.exit_code == 5
    assert "not found" in result.output.lower() or "no such file" in result.output.lower()
    assert "environment: ok" in result.output.lower()
    assert "markdown:" in result.output.lower()


def test_doctor_with_missing_input_supports_json_output(tmp_path: Path) -> None:
    result = runner.invoke(
        app,
        [
            "doctor",
            str(tmp_path / "missing.md"),
            "--style",
            "tests/fixtures/style.yaml",
            "--json",
        ],
    )

    assert result.exit_code == 5
    payload = json.loads(result.stdout)
    assert payload["success"] is False
    assert payload["summary"] == "Doctor failed"
    assert payload["checks"] == [
        {"name": "environment", "status": "ok"},
        {"name": "markdown", "status": "error"},
    ]
    assert payload["warnings"] == []
    assert "not found" in payload["error"]["message"].lower() or "no such file" in payload["error"]["message"].lower()


def test_doctor_returns_exit_code_4_for_missing_required_heading_style(
    tmp_path: Path,
) -> None:
    input_path = tmp_path / "heading3-only.md"
    input_path.write_text("### Third level title\n", encoding="utf-8")

    style_path = tmp_path / "paragraph-only.yaml"
    style_path.write_text(
        "\n".join(
            [
                "defaults:",
                "  font_name: Times New Roman",
                "elements:",
                "  paragraph:",
                "    based_on: Normal",
                "    font_size: 12",
            ]
        ),
        encoding="utf-8",
    )

    result = runner.invoke(
        app,
        [
            "doctor",
            str(input_path),
            "--style",
            str(style_path),
        ],
    )

    assert result.exit_code == 4
    assert "environment: ok" in result.output.lower()
    assert "markdown: ok" in result.output.lower()
    assert "style: ok" in result.output.lower()
    assert "heading3" in result.output.lower()
    assert "style:" in result.output.lower()


def test_doctor_with_style_resolution_failure_supports_json_output(
    tmp_path: Path,
) -> None:
    input_path = tmp_path / "heading3-only.md"
    input_path.write_text("### Third level title\n", encoding="utf-8")

    style_path = tmp_path / "paragraph-only.yaml"
    style_path.write_text(
        "\n".join(
            [
                "defaults:",
                "  font_name: Times New Roman",
                "elements:",
                "  paragraph:",
                "    based_on: Normal",
                "    font_size: 12",
            ]
        ),
        encoding="utf-8",
    )

    result = runner.invoke(
        app,
        [
            "doctor",
            str(input_path),
            "--style",
            str(style_path),
            "--json",
        ],
    )

    assert result.exit_code == 4
    payload = json.loads(result.stdout)
    assert payload["success"] is False
    assert payload["checks"] == [
        {"name": "environment", "status": "ok"},
        {"name": "markdown", "status": "ok", "detail": input_path.as_posix()},
        {"name": "style", "status": "ok", "detail": style_path.as_posix()},
        {"name": "style-resolution", "status": "error"},
    ]
    assert "heading3" in payload["error"]["message"].lower()


def test_doctor_reports_non_fatal_markdown_warnings(tmp_path: Path) -> None:
    input_path = tmp_path / "warn.md"
    input_path.write_text(
        "\n".join(
            [
                "# Title",
                "",
                "```py",
                "print(1)",
                "```",
                "",
                "![diagram](diagram.png)",
                "",
                "<div>raw html</div>",
            ]
        ),
        encoding="utf-8",
    )

    result = runner.invoke(
        app,
        [
            "doctor",
            str(input_path),
            "--style",
            "tests/fixtures/style.yaml",
        ],
    )

    assert result.exit_code == 0
    assert "warning:" in result.stdout.lower()
    assert "warnings: 3" in result.stdout.lower()
    assert "code block" in result.stdout.lower()
    assert "image" in result.stdout.lower()
    assert "html" in result.stdout.lower()


def test_doctor_json_includes_warning_summary(tmp_path: Path) -> None:
    input_path = tmp_path / "warn-json.md"
    input_path.write_text(
        "\n".join(
            [
                "# Title",
                "",
                "```py",
                "print(1)",
                "```",
                "",
                "![diagram](diagram.png)",
            ]
        ),
        encoding="utf-8",
    )

    result = runner.invoke(
        app,
        [
            "doctor",
            str(input_path),
            "--style",
            "tests/fixtures/style.yaml",
            "--json",
        ],
    )

    assert result.exit_code == 0
    payload = json.loads(result.stdout)
    assert payload["success"] is True
    assert len(payload["warnings"]) == 2
    assert any("code block" in warning for warning in payload["warnings"])
    assert any("images" in warning for warning in payload["warnings"])


def test_doctor_deduplicates_warning_categories(tmp_path: Path) -> None:
    input_path = tmp_path / "dup-warn.md"
    input_path.write_text(
        "\n".join(
            [
                "# Title",
                "",
                "```py",
                "print(1)",
                "```",
                "",
                "```js",
                "console.log(1)",
                "```",
                "",
                "![a](a.png)",
                "![b](b.png)",
            ]
        ),
        encoding="utf-8",
    )

    result = runner.invoke(
        app,
        [
            "doctor",
            str(input_path),
            "--style",
            "tests/fixtures/style.yaml",
        ],
    )

    assert result.exit_code == 0
    assert "warnings: 2" in result.stdout.lower()
    assert result.stdout.lower().count("warning:") == 2


def test_doctor_returns_exit_code_4_for_missing_template_file(tmp_path: Path) -> None:
    input_path = tmp_path / "doctor-template.md"
    input_path.write_text("# Title\n", encoding="utf-8")
    missing_template = tmp_path / "missing-template.docx"

    result = runner.invoke(
        app,
        [
            "doctor",
            str(input_path),
            "--style",
            "tests/fixtures/style.yaml",
            "--template",
            str(missing_template),
        ],
    )

    assert result.exit_code == 4
    assert "template" in result.output.lower()
    assert "style-resolution" in result.output.lower()


def test_doctor_returns_exit_code_4_for_missing_template_numbering_resource(tmp_path: Path) -> None:
    input_path = tmp_path / "ordered.md"
    input_path.write_text("1. First\n", encoding="utf-8")
    style_path = tmp_path / "ordered-style.yaml"
    style_path.write_text(
        "\n".join(
            [
                "defaults:",
                "  font_size: 12",
                "elements:",
                "  paragraph:",
                "    font_size: 12",
                "  ordered_list:",
                "    template_style: Normal",
                "    numbering_style: Normal",
            ]
        ),
        encoding="utf-8",
    )
    template_path = tmp_path / "template.docx"
    _write_template_docx(template_path)

    result = runner.invoke(
        app,
        [
            "doctor",
            str(input_path),
            "--style",
            str(style_path),
            "--template",
            str(template_path),
        ],
    )

    assert result.exit_code == 4
    assert "numbering" in result.output.lower()
    assert "style-resolution" in result.output.lower()


def test_doctor_returns_exit_code_4_for_missing_unordered_template_numbering_resource(tmp_path: Path) -> None:
    input_path = tmp_path / "unordered.md"
    input_path.write_text("- First\n", encoding="utf-8")
    style_path = tmp_path / "unordered-style.yaml"
    style_path.write_text(
        "\n".join(
            [
                "elements:",
                "  paragraph:",
                "    template_style: Normal",
                "  unordered_list:",
                "    template_style: Normal",
                "    numbering_style: Normal",
            ]
        ),
        encoding="utf-8",
    )
    template_path = tmp_path / "template.docx"
    Document().save(template_path)

    result = runner.invoke(
        app,
        [
            "doctor",
            str(input_path),
            "--style",
            str(style_path),
            "--template",
            str(template_path),
        ],
    )

    assert result.exit_code == 4
    assert "numbering" in result.output.lower()
    assert "style-resolution" in result.output.lower()


def test_doctor_reports_native_numbering_detail_when_template_resource_exists(tmp_path: Path) -> None:
    input_path = tmp_path / "ordered-ok.md"
    input_path.write_text("1. First\n", encoding="utf-8")
    style_path = tmp_path / "ordered-style.yaml"
    style_path.write_text(
        "\n".join(
            [
                "defaults:",
                "  font_size: 12",
                "elements:",
                "  paragraph:",
                "    font_size: 12",
                "  ordered_list:",
                "    template_style: List Number",
            ]
        ),
        encoding="utf-8",
    )
    template_path = tmp_path / "template.docx"
    _write_template_docx(template_path)

    result = runner.invoke(
        app,
        [
            "doctor",
            str(input_path),
            "--style",
            str(style_path),
            "--template",
            str(template_path),
            "--json",
        ],
    )

    assert result.exit_code == 0
    payload = json.loads(result.stdout)
    style_resolution = next(check for check in payload["checks"] if check["name"] == "style-resolution")
    assert style_resolution["status"] == "ok"
    assert "native numbering" in style_resolution["detail"].lower()
    assert "list number" in style_resolution["detail"].lower()


def test_style_explain_prints_resolved_style_json() -> None:
    result = runner.invoke(
        app,
        [
            "style",
            "explain",
            "tests/fixtures/style.yaml",
            "heading1",
        ],
    )

    assert result.exit_code == 0
    payload = json.loads(result.stdout)
    assert payload["element_name"] == "heading1"
    assert payload["word_style_name"] == "Heading 1"
    assert payload["resolved_properties"]["font_size"] == 18
    assert payload["source_map"]["font_size"] == "yaml"


def test_style_explain_supports_heading_fallback_resolution() -> None:
    result = runner.invoke(
        app,
        [
            "style",
            "explain",
            "tests/fixtures/style.yaml",
            "heading3",
        ],
    )

    assert result.exit_code == 0
    payload = json.loads(result.stdout)
    assert payload["element_name"] == "heading3"
    assert payload["word_style_name"] == "Heading 3"
    assert payload["resolved_properties"]["font_size"] == 18
    assert payload["source_map"]["font_size"] == "yaml"


def test_style_explain_supports_pretty_output() -> None:
    result = runner.invoke(
        app,
        [
            "style",
            "explain",
            "tests/fixtures/style.yaml",
            "heading1",
            "--pretty",
        ],
    )

    assert result.exit_code == 0
    assert result.stdout.startswith("{\n")
    assert '\n  "element_name"' in result.stdout
    payload = json.loads(result.stdout)
    assert payload["element_name"] == "heading1"


def test_style_explain_supports_template_backed_resolution(tmp_path: Path) -> None:
    style_path = tmp_path / "template-style.yaml"
    style_path.write_text(
        "\n".join(
            [
                "defaults:",
                "  font_name: Times New Roman",
                "elements:",
                "  heading1:",
                "    template_style: Heading 1",
                "    font_size: 18",
            ]
        ),
        encoding="utf-8",
    )
    template_path = tmp_path / "template.docx"
    _write_template_docx(template_path)

    result = runner.invoke(
        app,
        [
            "style",
            "explain",
            str(style_path),
            "heading1",
            "--template",
            str(template_path),
        ],
    )

    assert result.exit_code == 0
    payload = json.loads(result.stdout)
    assert payload["word_style_name"] == "Heading 1"
    assert payload["resolved_properties"]["font_name"] == "Courier New"
    assert payload["source_map"]["font_name"] == "template_style"


def test_style_explain_outputs_font_slots_paragraph_properties_and_numbering(tmp_path: Path) -> None:
    style_path = tmp_path / "template-style.yaml"
    style_path.write_text(
        "\n".join(
            [
                "defaults:",
                "  font_size: 12",
                "elements:",
                "  ordered_list:",
                "    template_style: List Number",
                "    font_east_asia: 仿宋",
            ]
        ),
        encoding="utf-8",
    )
    template_path = tmp_path / "template.docx"
    document = Document()
    list_style = document.styles["List Number"]
    list_style.font.name = "Times New Roman"
    list_style._element.get_or_add_rPr().get_or_add_rFonts().set(qn("w:eastAsia"), "宋体")
    list_style.paragraph_format.first_line_indent = Pt(21)
    document.add_paragraph("Seed", style="List Number")
    document.save(template_path)

    result = runner.invoke(
        app,
        [
            "style",
            "explain",
            str(style_path),
            "ordered_list",
            "--template",
            str(template_path),
        ],
    )

    assert result.exit_code == 0
    payload = json.loads(result.stdout)
    assert payload["word_style_name"] == "List Number"
    assert payload["resolved_properties"]["font_ascii"] == "Times New Roman"
    assert payload["resolved_properties"]["font_east_asia"] == "仿宋"
    assert payload["resolved_properties"]["first_line_indent_pt"] == 21
    assert payload["resolved_properties"]["numbering_style"] == "List Number"
    assert payload["source_map"]["font_ascii"] == "template_style"
    assert payload["source_map"]["font_east_asia"] == "yaml"
    assert payload["numbering"]["requested_style"] == "List Number"
    assert payload["numbering"]["source"] == "template_style"
    assert payload["numbering"]["available_in_template"] is True
    assert payload["numbering"]["fallback_mode"] == "native"


def test_style_explain_reports_text_prefix_fallback_when_numbering_resource_is_missing(tmp_path: Path) -> None:
    style_path = tmp_path / "template-style.yaml"
    style_path.write_text(
        "\n".join(
            [
                "defaults:",
                "  font_size: 12",
                "elements:",
                "  ordered_list:",
                "    template_style: Normal",
                "    numbering_style: Normal",
            ]
        ),
        encoding="utf-8",
    )
    template_path = tmp_path / "template.docx"
    _write_template_docx(template_path)

    result = runner.invoke(
        app,
        [
            "style",
            "explain",
            str(style_path),
            "ordered_list",
            "--template",
            str(template_path),
        ],
    )

    assert result.exit_code == 0
    payload = json.loads(result.stdout)
    assert payload["numbering"]["requested_style"] == "Normal"
    assert payload["numbering"]["available_in_template"] is False
    assert payload["numbering"]["fallback_mode"] == "text-prefix"


def test_style_merge_applies_overrides_and_prints_json() -> None:
    result = runner.invoke(
        app,
        [
            "style",
            "merge",
            "tests/fixtures/style.yaml",
            "--set",
            "document.toc.depth=4",
        ],
    )

    assert result.exit_code == 0
    payload = json.loads(result.stdout)
    assert payload["document"]["toc"]["depth"] == 4


def test_style_merge_supports_pretty_output() -> None:
    result = runner.invoke(
        app,
        [
            "style",
            "merge",
            "tests/fixtures/style.yaml",
            "--set",
            "document.toc.depth=4",
            "--pretty",
        ],
    )

    assert result.exit_code == 0
    assert result.stdout.startswith("{\n")
    assert '\n  "document"' in result.stdout
    payload = json.loads(result.stdout)
    assert payload["document"]["toc"]["depth"] == 4


def test_style_show_prints_builtin_style_yaml() -> None:
    result = runner.invoke(app, ["style", "show", "academic"])

    assert result.exit_code == 0
    assert "name: academic" in result.stdout
    assert "font_ascii: Times New Roman" in result.stdout


def test_style_help_lists_subcommands() -> None:
    result = runner.invoke(app, ["style", "--help"])

    assert result.exit_code == 0
    assert "Inspect and validate style profiles." in result.stdout
    assert "show" in result.stdout
    assert "validate" in result.stdout
    assert "explain" in result.stdout
    assert "merge" in result.stdout


def test_style_show_supports_json_output() -> None:
    result = runner.invoke(app, ["style", "show", "academic", "--json"])

    assert result.exit_code == 0
    payload = json.loads(result.stdout)
    assert payload["name"] == "academic"
    assert payload["path"] == "academic.yaml"
    assert "name: academic" in payload["content"]
    assert "font_ascii: Times New Roman" in payload["content"]


def test_style_validate_accepts_valid_style_file() -> None:
    result = runner.invoke(app, ["style", "validate", "tests/fixtures/style.yaml"])

    assert result.exit_code == 0
    assert "Style OK" in result.stdout


def test_style_validate_supports_json_output() -> None:
    result = runner.invoke(
        app,
        ["style", "validate", "tests/fixtures/style.yaml", "--json"],
    )

    assert result.exit_code == 0
    payload = json.loads(result.stdout)
    assert payload == {
        "success": True,
        "style_path": "tests/fixtures/style.yaml",
        "message": "Style OK",
    }


def test_style_validate_rejects_invalid_style_file(tmp_path: Path) -> None:
    bad_style = tmp_path / "bad-style.yaml"
    bad_style.write_text("defaults: []\n", encoding="utf-8")

    result = runner.invoke(app, ["style", "validate", str(bad_style)])

    assert result.exit_code == 4
    assert "dictionary" in result.output.lower()


def test_style_validate_invalid_file_supports_json_output(tmp_path: Path) -> None:
    bad_style = tmp_path / "bad-style.yaml"
    bad_style.write_text("defaults: []\n", encoding="utf-8")

    result = runner.invoke(app, ["style", "validate", str(bad_style), "--json"])

    assert result.exit_code == 4
    payload = json.loads(result.stdout)
    assert payload["success"] is False
    assert payload["style_path"] == bad_style.as_posix()
    assert "dictionary" in payload["message"].lower()
