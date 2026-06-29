from __future__ import annotations

from pathlib import Path
import base64

from pptx import Presentation
from pptx.util import Inches

from docubridge.core.pptx_ingest import parse_pptx_file

_PNG_1X1 = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+aRX0AAAAASUVORK5CYII="
)


def _write_sample_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Intro"
    slide.placeholders[1].text = "First point\nSecond point"
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)


def _write_image_pptx(path: Path, image_path: Path) -> None:
    image_path.parent.mkdir(parents=True, exist_ok=True)
    image_path.write_bytes(_PNG_1X1)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Images"
    picture = slide.shapes.add_picture(str(image_path), Inches(1), Inches(1))
    picture.name = "deck-image"
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)


def _write_notes_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Notes"
    slide.placeholders[1].text = "Visible point"
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.text = "Speaker note one"
    notes_frame.add_paragraph().text = "Speaker note two"
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)


def _write_table_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Table"
    table_shape = slide.shapes.add_table(3, 2, Inches(1), Inches(1.5), Inches(4), Inches(1.5))
    table = table_shape.table
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Alpha"
    table.cell(1, 1).text = "1"
    table.cell(2, 0).text = "Beta"
    table.cell(2, 1).text = "2"
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)


def _write_rich_table_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Table"
    table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(1.5), Inches(5), Inches(1.2))
    table = table_shape.table
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Value"

    first_cell = table.cell(1, 0)
    first_paragraph = first_cell.text_frame.paragraphs[0]
    first_paragraph.text = ""
    bold_run = first_paragraph.add_run()
    bold_run.text = "Alpha"
    bold_run.font.bold = True

    second_cell = table.cell(1, 1)
    second_paragraph = second_cell.text_frame.paragraphs[0]
    second_paragraph.text = ""
    second_paragraph.add_run().text = "See "
    link_run = second_paragraph.add_run()
    link_run.text = "ref"
    link_run.hyperlink.address = "https://example.com"

    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)


def _write_ordered_pptx(path: Path, image_path: Path) -> None:
    image_path.parent.mkdir(parents=True, exist_ok=True)
    image_path.write_bytes(_PNG_1X1)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Ordered"
    table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(3.4), Inches(4), Inches(1.2))
    table = table_shape.table
    table.cell(0, 0).text = "Key"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Order"
    table.cell(1, 1).text = "Preserved"
    picture = slide.shapes.add_picture(str(image_path), Inches(1), Inches(2.2))
    picture.name = "ordered-image"
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(4), Inches(0.6))
    textbox.text_frame.text = "Before image"
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)


def _write_text_hierarchy_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Hierarchy"
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(5), Inches(2))
    text_frame = textbox.text_frame
    text_frame.text = "Overview paragraph"
    first_item = text_frame.add_paragraph()
    first_item.text = "Top item"
    first_item.level = 0
    nested_item = text_frame.add_paragraph()
    nested_item.text = "Nested item"
    nested_item.level = 1
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)


def _write_rich_text_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Rich Text"
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.4), Inches(6), Inches(1.5))
    text_frame = textbox.text_frame
    paragraph = text_frame.paragraphs[0]
    paragraph.text = ""
    paragraph.add_run().text = "Body "
    bold_run = paragraph.add_run()
    bold_run.text = "Bold"
    bold_run.font.bold = True
    plain_run = paragraph.add_run()
    plain_run.text = " "
    italic_run = paragraph.add_run()
    italic_run.text = "Italic"
    italic_run.font.italic = True
    plain_run_two = paragraph.add_run()
    plain_run_two.text = " "
    link_run = paragraph.add_run()
    link_run.text = "ref"
    link_run.hyperlink.address = "https://example.com"
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)


def _write_code_text_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Code"

    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.4), Inches(6), Inches(1.2))
    text_frame = textbox.text_frame
    paragraph = text_frame.paragraphs[0]
    paragraph.text = ""
    paragraph.add_run().text = "Body "
    code_run = paragraph.add_run()
    code_run.text = "code"
    code_run.font.name = "Consolas"

    table_shape = slide.shapes.add_table(2, 1, Inches(1), Inches(3), Inches(4), Inches(1.2))
    table = table_shape.table
    table.cell(0, 0).text = "Value"
    cell_paragraph = table.cell(1, 0).text_frame.paragraphs[0]
    cell_paragraph.text = ""
    cell_paragraph.add_run().text = "Cell "
    cell_code = cell_paragraph.add_run()
    cell_code.text = "code"
    cell_code.font.name = "Consolas"

    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)


def test_parse_pptx_file_extracts_slide_title_and_body(tmp_path: Path) -> None:
    input_path = tmp_path / "sample.pptx"
    _write_sample_pptx(input_path)

    content = parse_pptx_file(input_path)

    assert "### Slide 1" in content
    assert "#### Intro" in content
    assert "- First point" in content
    assert "- Second point" in content


def test_parse_pptx_file_exports_images_as_markdown_references(tmp_path: Path) -> None:
    input_path = tmp_path / "images.pptx"
    image_path = tmp_path / "diagram.png"
    media_dir = tmp_path / "media"
    _write_image_pptx(input_path, image_path)

    content = parse_pptx_file(input_path, media_dir=media_dir)

    assert "#### Images" in content
    assert "![deck-image](" in content
    assert "assets/" in content
    exported = list(media_dir.iterdir())
    assert exported


def test_parse_pptx_file_extracts_speaker_notes_as_blockquotes(tmp_path: Path) -> None:
    input_path = tmp_path / "notes.pptx"
    _write_notes_pptx(input_path)

    content = parse_pptx_file(input_path)

    assert "#### Notes" in content
    assert "- Visible point" in content
    assert "> Notes:" in content
    assert "> Speaker note one" in content
    assert "> Speaker note two" in content


def test_parse_pptx_file_extracts_tables_as_markdown(tmp_path: Path) -> None:
    input_path = tmp_path / "table.pptx"
    _write_table_pptx(input_path)

    content = parse_pptx_file(input_path)

    assert "#### Table" in content
    assert "| Name | Value |" in content
    assert "| --- | --- |" in content
    assert "| Alpha | 1 |" in content
    assert "| Beta | 2 |" in content


def test_parse_pptx_file_preserves_rich_text_inside_table_cells(tmp_path: Path) -> None:
    input_path = tmp_path / "rich-table.pptx"
    _write_rich_table_pptx(input_path)

    content = parse_pptx_file(input_path)

    assert "#### Table" in content
    assert "| **Alpha** | See [ref](https://example.com) |" in content


def test_parse_pptx_file_preserves_block_order_for_text_images_and_tables(tmp_path: Path) -> None:
    input_path = tmp_path / "ordered.pptx"
    image_path = tmp_path / "ordered.png"
    media_dir = tmp_path / "media"
    _write_ordered_pptx(input_path, image_path)

    content = parse_pptx_file(input_path, media_dir=media_dir)

    text_index = content.index("Before image")
    image_index = content.index("![ordered-image](")
    table_index = content.index("| Key | Value |")
    assert text_index < image_index < table_index


def test_parse_pptx_file_distinguishes_plain_text_and_nested_list_items(tmp_path: Path) -> None:
    input_path = tmp_path / "hierarchy.pptx"
    _write_text_hierarchy_pptx(input_path)

    content = parse_pptx_file(input_path)

    assert "Overview paragraph" in content
    assert "- Top item" in content
    assert "  - Nested item" in content


def test_parse_pptx_file_preserves_rich_text_in_textboxes(tmp_path: Path) -> None:
    input_path = tmp_path / "rich-text.pptx"
    _write_rich_text_pptx(input_path)

    content = parse_pptx_file(input_path)

    assert "#### Rich Text" in content
    assert "Body **Bold** *Italic* [ref](https://example.com)" in content


def test_parse_pptx_file_preserves_inline_code_in_textboxes_and_tables(tmp_path: Path) -> None:
    input_path = tmp_path / "code-text.pptx"
    _write_code_text_pptx(input_path)

    content = parse_pptx_file(input_path)

    assert "#### Code" in content
    assert "Body `code`" in content
    assert "| Cell `code` |" in content


def test_parse_pptx_file_preserves_inline_code_in_textboxes_and_table_cells(tmp_path: Path) -> None:
    input_path = tmp_path / "code-text.pptx"
    _write_code_text_pptx(input_path)

    content = parse_pptx_file(input_path)

    assert "#### Code" in content
    assert "Body `code`" in content
    assert "| Cell `code` |" in content
