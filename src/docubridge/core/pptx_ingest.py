from __future__ import annotations

from pathlib import Path
import re

from pptx import Presentation


def _slugify_name(value: str) -> str:
    text = re.sub(r"[^A-Za-z0-9._-]+", "-", value.strip()).strip("-")
    return text or "image"


def _export_shape_image(shape, media_dir: Path, image_index: int) -> tuple[str, str]:
    image = shape.image
    alt_text = _slugify_name(Path(getattr(shape, "name", "") or "").stem)
    extension = image.ext or "png"
    filename = f"{image_index:03d}-{alt_text}.{extension}"
    media_dir.mkdir(parents=True, exist_ok=True)
    output_path = media_dir / filename
    output_path.write_bytes(image.blob)
    return alt_text, f"assets/{filename}"


def _extract_notes_lines(slide) -> list[str]:
    try:
        text_frame = slide.notes_slide.notes_text_frame
    except AttributeError:
        return []

    lines: list[str] = []
    for paragraph in text_frame.paragraphs:
        text = paragraph.text.strip()
        if text:
            lines.append(text)
    return lines


def _extract_text_frame_lines(text_frame) -> list[str]:
    lines = []
    for paragraph in text_frame.paragraphs:
        text = _extract_paragraph_text(paragraph)
        if text:
            lines.append(text)
    return lines


def _extract_table_lines(shape) -> list[str]:
    rows = []
    for row in shape.table.rows:
        cells = ["<br>".join(_extract_text_frame_lines(cell.text_frame)) for cell in row.cells]
        if any(cells):
            rows.append(cells)
    if not rows:
        return []

    header = rows[0]
    body = rows[1:]
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join("---" for _ in header) + " |",
    ]
    for row in body:
        padded = row + [""] * (len(header) - len(row))
        lines.append("| " + " | ".join(padded[: len(header)]) + " |")
    return lines


def _shape_position(shape) -> tuple[int, int]:
    top = int(getattr(shape, "top", 0) or 0)
    left = int(getattr(shape, "left", 0) or 0)
    return (top, left)


def _wrap_markdown_text(
    text: str,
    *,
    bold: bool = False,
    italic: bool = False,
) -> str:
    if not text:
        return ""
    if bold and italic:
        return f"***{text}***"
    if bold:
        return f"**{text}**"
    if italic:
        return f"*{text}*"
    return text


def _is_code_run(run) -> bool:
    font_name = (getattr(run.font, "name", "") or "").strip().lower()
    return font_name in {"consolas", "courier new", "courier"}


def _extract_paragraph_text(paragraph) -> str:
    parts: list[str] = []
    for run in paragraph.runs:
        text = _wrap_markdown_text(
            run.text,
            bold=bool(getattr(run.font, "bold", False)),
            italic=bool(getattr(run.font, "italic", False)),
        )
        if text and _is_code_run(run):
            text = f"`{text}`"
        href = getattr(run.hyperlink, "address", "") if getattr(run, "hyperlink", None) is not None else ""
        if href and text:
            parts.append(f"[{text}]({href})")
            continue
        parts.append(text)
    return "".join(parts).strip()


def _extract_text_lines(shape) -> list[str]:
    if not hasattr(shape, "text_frame"):
        return []

    paragraphs = []
    for paragraph in shape.text_frame.paragraphs:
        text = _extract_paragraph_text(paragraph)
        if text:
            paragraphs.append((text, int(getattr(paragraph, "level", 0) or 0)))
    if not paragraphs:
        return []

    is_placeholder = bool(getattr(shape, "is_placeholder", False))
    if is_placeholder:
        lines: list[str] = []
        for text, level in paragraphs:
            for raw_line in text.splitlines():
                line = raw_line.strip()
                if not line:
                    continue
                indent = "  " * max(level, 0)
                lines.append(f"{indent}- {line}")
        return lines

    if len(paragraphs) == 1:
        return [paragraphs[0][0]]

    lines: list[str] = [paragraphs[0][0]]
    for text, level in paragraphs[1:]:
        indent = "  " * max(level, 0)
        lines.append(f"{indent}- {text}")
    return lines


def parse_pptx_file(path: Path, media_dir: Path | None = None) -> str:
    presentation = Presentation(path)
    lines: list[str] = []
    image_index = 1

    for index, slide in enumerate(presentation.slides, start=1):
        lines.append(f"### Slide {index}")
        lines.append("")

        title = slide.shapes.title.text.strip() if slide.shapes.title and slide.shapes.title.text else ""
        if title:
            lines.append(f"#### {title}")
            lines.append("")

        content_shapes = [shape for shape in slide.shapes if shape != slide.shapes.title]
        for shape in sorted(content_shapes, key=_shape_position):
            if hasattr(shape, "table"):
                lines.extend(_extract_table_lines(shape))
                lines.append("")
                continue
            if hasattr(shape, "image") and media_dir is not None:
                alt_text, relative_path = _export_shape_image(shape, media_dir, image_index)
                lines.append(f"![{alt_text}]({relative_path})")
                lines.append("")
                image_index += 1
                continue
            text_lines = _extract_text_lines(shape)
            if not text_lines:
                continue
            lines.extend(text_lines)

        notes_lines = _extract_notes_lines(slide)
        if notes_lines:
            lines.append("")
            lines.append("> Notes:")
            for note in notes_lines:
                lines.append(f"> {note}")
        lines.append("")

    while lines and lines[-1] == "":
        lines.pop()
    return "\n".join(lines) + ("\n" if lines else "")
